"""
This script used the json file from the Handballfreunde Statistik App.
Execute with python handball_stats_visualisation.py.
Analysis can then be found in the output folder.
"""

import os
import json
from os import makedirs
from os.path import isdir

import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

from constants import missed_shots_list, CURRENT_DIR, PATH, OUTPUT_DIR, PPT_FILE_PATH, PPT_FILE_PATH_NEW
from ppt_creation import create_ppt
from shot_analysis import analyze_keeper

ppt = Presentation()

# TODO location translation ("K" -> "Kreis")
# TODO time per attack ist eine analyse, das hat mit parsing wenig zu tun
def parse_json(json_string: str) -> pd.DataFrame:
    raw_data = json.loads(json_string)
    data = pd.DataFrame(raw_data['actions'])

    # sort for time
    data.sort_values('timestamp', inplace=True)

    # handle None values
    data['player'] = ["player undefined" if i is None else i for i in data['player']]
    data['x'] = ['0' if i is None else i for i in data['x']]
    data['y'] = ['0' if i is None else i for i in data['y']]
    data['x'] = data['x'].astype(int)
    data['y'] = data['y'].astype(int)

    # split up into minute and second
    data['minute'] = round(data['timestamp'] / 60, 0).astype(int)
    data['seconds'] = (data['timestamp'] % 60).astype(int)

    # make extra column for nr, name and position
    # FIXME: use existing column and filter dict directly with sth like below
    # formatted_data[formatted_data['player'].apply(lambda x: x['player_name'] == 'Ben Loius Vollert')
    for i in range(0, len(data)):
        if data.loc[i, 'player'] != 'player undefined':
            if not pd.isna(data.loc[i, 'player']['number']):
                data.loc[i, 'number'] = int(data.loc[i, 'player']['number'])
            if not pd.isna(data.loc[i, 'player']['player_name']):
                data.loc[i, 'player_name'] = data.loc[i, 'player']['player_name']
        if i == 0:
            data.loc[i, 'attack_time'] = data.loc[i, 'timestamp']
        else:
            data.loc[i, 'attack_time'] = data.loc[i, 'timestamp'] - data.loc[i - 1, 'timestamp']

    data = data.drop('player', axis=1)
    data.fillna(0, inplace=True)
    data['number'] = data['number'].astype(int)
    data['attack_time'] = data['attack_time'].astype(int)

    return data


def full_game_analysis(formatted_data):
    """Analyse over 60 minutes when we scored, missed or made a mistake."""

    # game analyses: Treffer, Fehlwurf, Fehlpass, Teschnischer Fehler
    # FIXME: define these dfs in separate funciton
    df_shots = formatted_data[formatted_data['own_team']]
    df_score = df_shots[~(df_shots['type'].isin(missed_shots_list))]
    df_fehler = df_shots[df_shots['type'] == 'Fehler']
    df_miss = df_shots[(df_shots['type'].isin(missed_shots_list))]

    spiel_df = pd.DataFrame(list(range(61)), columns=['Minute im Spiel'])
    for i in range(0, len(spiel_df)):
        treffer_list = df_score["minute"].value_counts()[1:]
        if i in treffer_list:
            spiel_df.loc[i, "Treffer"] = treffer_list[i]

        fehler_list = df_fehler["minute"].value_counts()[1:]
        if i in fehler_list:
            spiel_df.loc[i, "Fehler"] = fehler_list[i]

        verworfen_list = df_miss["minute"].value_counts()[1:]
        if i in verworfen_list:
            spiel_df.loc[i, "Verworfen"] = verworfen_list[i]

    count = 0
    for column in spiel_df.columns[1:]:
        if column == 'Treffer':
            color = 'green'
        elif column == 'Verworfen':
            color = 'orange'
        else:
            color = 'red'

        # plot each column with an offset by shifting the index
        # FIXME: write function for plots and ppt
        plt.figure(figsize=(20, 2.5))
        plt.xlim(0, 61)
        plt.xticks(spiel_df.index)
        plt.bar(spiel_df.index, spiel_df[column], width=0.5, color=color, label=column)
        plt.legend(loc='upper left')

        img_path = os.path.join(OUTPUT_DIR, f"plot_spiel_{column}.png")
        plt.savefig(img_path)
        plt.close()

        # add all graphs to one slide
        slide_layout = ppt.slide_layouts[5]
        if count == 0:
            slide = ppt.slides.add_slide(slide_layout)
        move_down = 1 + 2 * count
        slide.shapes.add_picture(img_path, Inches(-1), Inches(move_down), width=Inches(11.75), height=Inches(1.75))
        count += 1


def seconds_per_attack(formatted_data):
    """Visualise time in attack for each team."""

    # check time in attack for handballfreunde
    # FIXME: use ball changing teams and not time until shot for time in attack
    seconds_in_attack = formatted_data[formatted_data['own_team']]['attack_time']
    fig = plt.figure(figsize=(20, 5))  # noqa: F841, needed for format in ppt
    plt.bar(height=seconds_in_attack, x=list(formatted_data.loc[seconds_in_attack.index, 'minute']))
    plt.title('Handballfreunde Sekunden pro Angriff')

    # save image
    # FIXME: Make function for image and ppt generation for 60min graphs (also starting slide)
    img_path = os.path.join(OUTPUT_DIR, "plot_spiel_hbf_time_attack.png")
    plt.savefig(img_path)
    plt.close()

    # add all graphs to one slide
    """slide_layout = ppt.slide_layouts[5]
    slide = ppt.slides.add_slide(slide_layout)
    move_down = 1
    slide.shapes.add_picture(img_path, Inches(-1), Inches(move_down), width=Inches(11.75), height=Inches(3))"""

    # check time in attack for opponents
    opponent_seconds_in_attack = formatted_data[~formatted_data['own_team']]['attack_time']
    plt.figure(figsize=(20, 5))
    plt.bar(height=opponent_seconds_in_attack, x=list(formatted_data.loc[opponent_seconds_in_attack.index, 'minute']))
    plt.title('Gegner Sekunden pro Angriff')

    # save image
    img_path = os.path.join(OUTPUT_DIR, "plot_spiel_opponent_time_attack.png")
    plt.savefig(img_path)
    plt.close()

    # add all graphs to one slide
    move_down = 4
    #slide.shapes.add_picture(img_path, Inches(-1), Inches(move_down), width=Inches(11.75), height=Inches(3))


def opponent_analysis_table(formatted_data):
    """Analyse from which position the opponent scored."""
    # FIXME: Make this position_analysis and include table for handballfreunde
    # FIXME: Add 'Konter' as position

    # create oppenet df
    df_opponent_shots = formatted_data[~formatted_data['own_team']]
    df_opponents_scored = df_opponent_shots[~(df_opponent_shots['type'].isin(missed_shots_list))]['location'].value_counts().reset_index()
    df_opponents_missed = df_opponent_shots[(df_opponent_shots['type'].isin(missed_shots_list))]['location'].value_counts().reset_index()

    # differentiate per position shot and miss
    df_opponents_scored.rename(columns={"location": "Position (Gegner)", "count": "Treffer"}, inplace=True)
    df_opponents_missed.rename(columns={"location": "Position (Gegner)", "count": "Verworfen"}, inplace=True)
    opponent_analysis = df_opponents_scored.merge(df_opponents_missed, how='outer', on='Position (Gegner)')

    opponent_analysis.fillna(0, inplace=True)
    opponent_analysis.sort_values('Position (Gegner)', ascending=True)
    opponent_analysis.set_index('Position (Gegner)', inplace=True)

    # calculate quote per position
    opponent_analysis['Quote'] = round(
        opponent_analysis['Treffer'] / (opponent_analysis['Treffer'] + opponent_analysis['Verworfen']) * 100, 2)
    opponent_analysis = opponent_analysis.astype(int)
    opponent_analysis.reset_index(inplace=True)

    # save table
    img_path = os.path.join(OUTPUT_DIR, "opponent_analysis.png")
    df_to_image(opponent_analysis, img_path)

    # add to ppt
    #add_to_ppt(img_path, 1, 1, 8, 5)


def shot_visualisation(df_shots, team, player=None, location=None):
    """Visualisa shots."""
    # FIXME: Make heatmap instead of normal plot

    # check whether visualisation for attack or defense (keeper)
    # FIXME: Move this in extra function that outputs each df
    if team == 'handballfreunde':
        df_shots = df_shots[df_shots['own_team']]
        success_color = 'green'
        fail_color = 'red'
    else:
        df_shots = df_shots[~df_shots['own_team']]
        success_color = 'red'  # for keeper a missed shot is good
        fail_color = 'green'

    # check whether for entire team or specific player
    if player is None:
        df_score = df_shots[~(df_shots['type'].isin(missed_shots_list))]
        df_miss = df_shots[(df_shots['type'].isin(missed_shots_list))]
    else:
        # FIXME: Filter out 7M and visualise separatly
        df_score = df_shots[(df_shots['player_name'] == player) & ~(df_shots['type'].isin(missed_shots_list))]
        df_miss = df_shots[(df_shots['player_name'] == player) & (df_shots['type'].isin(missed_shots_list))]

    # calculate quote
    scored = len(df_score)
    missed = len(df_miss)
    shots = missed + scored
    try:
        if team == 'handballfreunde':
            quote = round((scored / (shots) * 100), 2)
        else:
            quote = round((missed / (shots) * 100), 2)
    except ZeroDivisionError:
        quote = "0"

    # plot one graph
    plt.figure()
    # always show entire goal (15x10)
    plt.xlim(0, 100)
    plt.xticks(())
    plt.ylim(0, 100)
    plt.yticks(())
    plt.plot(df_miss["x"], df_miss["y"], marker='o', linestyle='none', color=fail_color, ms=17,
             label='Verworfen')
    plt.plot(df_score["x"], df_score["y"], marker='o', linestyle='none', color=success_color, ms=15,
             label='Treffer')
    # set the title and labels
    if player is None:
        plt.title(
                f"Wurfanalyse für Handballfreunde:\nWürfe =  {shots}, Teffer = {scored}, Quote = {quote}%"
                )
        img_path = os.path.join(OUTPUT_DIR, "plot_team.png")
    else:
        if team == 'handballfreunde':
            plt.title(
                        f"Wurfanalyse für {player}:\nWürfe =  {shots}, Teffer = {scored}, Quote = {quote}%"
                        )
        else:
            if location is None:
                plt.title(
                            f"Keeperanalyse für {player}:\nWürfe =  {shots}, Gehalten = {missed}, Quote = {quote}%"
                            )
            else:
                plt.title(
                            f"Keeperanalyse für {player} von {location}:\nWürfe =  {shots}, Gehalten = {missed}, Quote = {quote}%"
                            )
        img_path = os.path.join(OUTPUT_DIR, f"plot_{player}.png")

    # save image
    plt.savefig(img_path)
    plt.close()

    # add to ppt
    add_to_ppt(img_path, 1, 1, 8, 5)


def add_to_ppt(img_path, left, top, width=None, height=None):
    # Add graph to ppt

    slide_layout = ppt.slide_layouts[5]
    slide = ppt.slides.add_slide(slide_layout)
    slide.shapes.add_picture(img_path, left=Inches(left), top=Inches(top), width=Inches(width), height=Inches(height))


def df_to_image(df, path):
    # Convert pandas dataframe to image to be able to paste it into ppt.
    fig, ax = plt.subplots(figsize=(8, 5))
    ax.axis('off')
    table = ax.table(cellText=df.values, colLabels=df.columns, cellLoc='center', loc='center')
    table.auto_set_font_size(False)
    table.set_fontsize(12)
    table.scale(1.2, 1.2)

    plt.savefig(path, bbox_inches='tight')
    plt.close(fig)

# TODO create input folder and parse every file in it
def main():

    if not isdir(OUTPUT_DIR):
        makedirs(OUTPUT_DIR)

    # add title slide
    title_img_path = rf"{CURRENT_DIR}/woran-hat-es-gelegen-winner.png"
    add_to_ppt(title_img_path, 2, 0, 6, 7)

    # Import data from json and parse it for further analysis.
    with open(PATH) as json_data:
        data = parse_json(json_data.read())
        json_data.close()

    # TODO struktur anpassen
    # ab hier dann analysen einbauen, denen "data" übergeben
    # die geben dann entweder bilder oder slides zurück
    # danach mergen wir dann die ganzen analysen (evtl. auch mit inhaltsverzeichnis?) (geht das auch nur mit pyPdf?)

    images = {"Keeper": analyze_keeper(data)}

    df_shots = data[data['type'] != 'Fehler']

    # visualise data and export to ppt
    full_game_analysis(data)
    seconds_per_attack(data)
    opponent_analysis_table(df_shots)
    shot_visualisation(df_shots, 'handballfreunde')

    # analysis per fieldplayer
    for player in df_shots[df_shots['own_team']]['player_name'].unique():
        shot_visualisation(df_shots, 'handballfreunde', player)

    ppt_new = create_ppt(images)

    # TODO platform based ppt/pdf saving
    # export to pdf
    try:
        ppt.save(PPT_FILE_PATH)
        ppt_new.save(PPT_FILE_PATH_NEW)
    except PermissionError:
        print("Kollege! Mach die PowerPoint zu, es zieht! Dann nochmal probieren bitte (-_-)")
        os.system("TASKKILL /F /IM powerpnt.exe")


if __name__ == '__main__':
    main()
