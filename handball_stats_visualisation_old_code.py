import pandas as pd
import matplotlib.pyplot as plt
from pandas import DataFrame
from pptx import Presentation
from pptx.util import Inches
import os
import sys
import comtypes.client

#FIXME: per player show minutes of score and misses
#FIXME: Add disclaimer via slide.addpicture() on title slide
#FIXME: Wurfanalyse Team auch auf 10x15 limitieren

def format_df(old_df: DataFrame, df_new: DataFrame) -> DataFrame:
    # data formatting
    count = 0
    for i in range(0, len(old_df)):
        nr = old_df.loc[i, 'Nr.']

        # split relevant columns at ;
        for column in old_df.columns[3:]:

            if old_df.loc[i, column][-1:] == ';':
                old_df.loc[i, column] = old_df.loc[i, column][:-1]
            entry_list = old_df.loc[i, column].split(';')

            for entry in entry_list:

                df_new.loc[count, 'Nr.'] = nr
                df_new.loc[count, f'{column}_min'] = entry.split(',')[0]
                if column in ['Treffer', 'Verworfen']:
                    df_new.loc[count, f'{column}_y'] = entry.split(',')[2]
                    df_new.loc[count, f'{column}_x'] = entry.split(',')[1]
                if column == 'Technisches Foul':
                    df_new.loc[count, f'{column}_position'] = entry.split(',')[1]
                count = count + 1

    df_new = df_new.fillna(0)

    return df_new


def PPTtoPDF(inputFileName, outputFileName, formatType=32):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    if outputFileName[-3:] != 'pdf':
        outputFileName = outputFileName + ".pdf"
    deck = powerpoint.Presentations.Open(inputFileName)
    deck.SaveAs(outputFileName, formatType)
    deck.Close()
    powerpoint.Quit()


def df_to_image(df, path):
    fig, ax = plt.subplots(figsize=(8, 5))
    ax.axis('off')
    table = ax.table(cellText=df.values, colLabels=df.columns, cellLoc='center', loc='center')
    table.auto_set_font_size(False)
    table.set_fontsize(12)
    table.scale(1.2, 1.2)

    plt.savefig(path, bbox_inches='tight')
    plt.close(fig)


def main():
    # TODO temporary
    print('Moin, wie heißt die Excel? (ohne .xlsx Endung)')
    filename = input()

    current_dir = os.getcwd()
    output_dir = current_dir + "/output/"
    os.makedirs(output_dir, exist_ok=True)
    input_file_path = rf"{current_dir}\{filename}.xlsx"

    print("Output to \n" + output_dir)

    # TODO refactoring zu einem dataframe
    df_tw, keeper_nr, nr, parsed_data, raw_data = parse_data(input_file_path)

    name_dict = {}
    for i in range(0, len(raw_data)):
        name_dict[raw_data.loc[i, 'Nr.'].astype(int)] = raw_data.loc[i, 'Spieler']

    # game analyses: Treffer, Fehlwurf, Fehlpass, Teschnischer Fehler
    df_shots = parsed_data
    spiel_df = pd.DataFrame(list(range(61)), columns=['Minute im Spiel'])
    for i in range(0, len(spiel_df)):
        treffer_list = df_shots["Treffer_min"].value_counts()[1:]
        if i in treffer_list:
            spiel_df.loc[i, "Treffer"] = treffer_list[i]

        fehlpass_list = df_shots["Fehlpass_min"].value_counts()[1:]
        if i in fehlpass_list:
            spiel_df.loc[i, "Fehlpass"] = fehlpass_list[i]

        verworfen_list = df_shots["Verworfen_min"].value_counts()[1:]
        if i in verworfen_list:
            spiel_df.loc[i, "Verworfen"] = verworfen_list[i]

        tf_list = df_shots["Technisches Foul_min"].value_counts()[1:]
        if i in tf_list:
            spiel_df.loc[i, "Technisches Foul"] = tf_list[i]


    # add game graphs to ppt
    ppt = Presentation()

    # titel slide
    title_img_path = rf"{current_dir}\woran-hat-es-gelegen-winner.png"

    # add to ppt
    slide_layout = ppt.slide_layouts[5]
    slide = ppt.slides.add_slide(slide_layout)
    left = Inches(2)
    top = Inches(1)
    slide.shapes.add_picture(title_img_path, left, top, width=Inches(6), height=Inches(6))

    count = 1
    for column in spiel_df.columns[1:]:

        if column == 'Treffer':
            color = 'green'
        if column == 'Verworfen':
            color = 'orange'
        if column == 'Fehlpass':
            color = 'red'
        if column == 'Technisches Foul':
            color = 'blue'

        # plot each column with an offset by shifting the index
        # FIXME: write function for plots and ppt
        plt.figure(figsize=(20, 2.5))
        plt.xlim(0, 61)
        plt.xticks(spiel_df.index)
        plt.bar(spiel_df.index, spiel_df[column], width=0.5, color=color, label=column)
        plt.legend(loc='upper left')

        # save image
        img_path = os.path.join(output_dir, f"plot_spiel_{column}.png")
        plt.savefig(img_path)
        plt.close()

        # add to ppt
        slide_layout = ppt.slide_layouts[5]
        if count == 1:
            slide = ppt.slides.add_slide(slide_layout)
        top = Inches(count + 1)
        left = Inches(0)
        slide.shapes.add_picture(img_path, left, top, width=Inches(10), height=Inches(1))

        count += 1


    # calculate quote for whole Team
    df_shots = parsed_data[~parsed_data['Nr.'].isin(keeper_nr)]
    treffer = len(df_shots[df_shots['Treffer_min'] != 0])
    gehalten = len(df_shots[df_shots['Verworfen_min'] != 0])
    quote = round((treffer / (treffer + gehalten) * 100), 2)
    # add figure for whole team
    plt.figure()
    plt.plot(df_shots["Verworfen_x"], df_shots["Verworfen_y"], marker='o', linestyle='none', color='red', ms=17)
    plt.plot(df_shots["Treffer_x"], df_shots["Treffer_y"], marker='o', linestyle='none', color='green', ms=15)
    plt.title(f"Wurfanalyse Handballfreunde: Quote = {quote} %")
    plt.xlabel('Torbreite 3m = 1-15')
    plt.ylabel('Torhöhe 2m = 1-10')
    # save image
    img_path = os.path.join(output_dir, f"plot_nr_{nr}.png")
    plt.savefig(img_path)
    plt.close()
    # add to ppt
    slide_layout = ppt.slide_layouts[5]
    slide = ppt.slides.add_slide(slide_layout)
    left = Inches(1)
    top = Inches(1)
    slide.shapes.add_picture(img_path, left, top, width=Inches(8), height=Inches(5))

    # gegner analyse
    df_opponent_shots = pd.DataFrame(
        columns=['Nr.', 'Treffer_min', 'Treffer_x', 'Treffer_y', 'Verworfen_min', 'Verworfen_x', 'Verworfen_y'])
    count = 0
    # loop for each keepr
    for i in range(0, len(df_tw)):
        nr = df_tw.loc[i, 'Nr.']
        # split relevant columns at ;
        # FIXME: Use format_df()
        # FIXME: unify with keeper analysis

        # loop for each column (Treffer, Verworfen, Fehlpass, Technisches Foul)
        for column in df_tw.columns[3:]:
            if df_tw.loc[i, column][-1:] == ';':
                df_tw.loc[i, column] = df_tw.loc[i, column][:-1]
            entry_list = df_tw.loc[i, column].split(';')
            # df_nr = df_new

            # loop for each event (RM,14,9 --> shot from RM in upper right corner)
            for entry in entry_list:
                df_opponent_shots.loc[count, 'Nr.'] = nr
                df_opponent_shots.loc[count, f'{column}_min'] = entry.split(',')[0]
                if column in ['Treffer', 'Verworfen']:
                    df_opponent_shots.loc[count, f'{column}_y'] = entry.split(',')[2]
                    df_opponent_shots.loc[count, f'{column}_x'] = entry.split(',')[1]
                if column == 'Technisches Foul':
                    df_opponent_shots.loc[count, f'{column}_position'] = entry.split(',')[1]
                count = count + 1

        df_opponent_shots = df_opponent_shots.fillna(0)

    gegner_analyse = df_opponent_shots['Treffer_min'][df_opponent_shots['Treffer_min'] != 0].value_counts().reset_index()
    gegner_analyse.rename(columns={"Treffer_min": "Position (Gegner)", "count": "Treffer"}, inplace=True)
    gegner_analyse.set_index('Position (Gegner)', inplace=True)
    gegner_verworfen = pd.DataFrame(df_opponent_shots['Verworfen_min'][df_opponent_shots['Verworfen_min'] != 0].value_counts())
    gegner_verworfen.reset_index(inplace=True)
    gegner_verworfen.rename(columns={"Verworfen_min": "Position (Gegner)", "count": "Verworfen"}, inplace=True)
    gegner_analyse.reset_index(inplace=True)
    gegner_analyse = gegner_analyse.merge(gegner_verworfen, how='outer')
    gegner_analyse.fillna(0, inplace=True)
    gegner_analyse.sort_values('Position (Gegner)', ascending=True)
    gegner_analyse.set_index('Position (Gegner)', inplace=True)
    gegner_analyse['Quote'] = round(
        gegner_analyse['Treffer'] / (gegner_analyse['Treffer'] + gegner_analyse['Verworfen']) * 100, 2)
    gegner_analyse = gegner_analyse.astype(int)
    gegner_analyse.reset_index(inplace=True)
    # add df to ppt
    img_path = os.path.join(output_dir, f"gegner_analyse.png")
    df_to_image(gegner_analyse, img_path)
    slide_layout = ppt.slide_layouts[5]
    slide = ppt.slides.add_slide(slide_layout)
    left = Inches(1)
    top = Inches(1)
    slide.shapes.add_picture(img_path, left, top, width=Inches(8), height=Inches(5))


    # shots in attack and saves in defense per player
    for nr in parsed_data['Nr.'].unique():

        # filter for nr
        df_nr = parsed_data[parsed_data['Nr.'] == int(nr)]
        df_nr = parsed_data[(parsed_data['Nr.'] == int(nr))]
        treffer = len(df_nr[df_nr['Treffer_min'] != 0])
        gehalten = len(df_nr[df_nr['Verworfen_min'] != 0])
        würfe = treffer + gehalten
        try:
            quote = round((treffer / (würfe) * 100), 2)
        except:
            quote = "0"

        # plot one graph per nr.
        plt.figure()
        plt.plot(df_nr["Verworfen_x"], df_nr["Verworfen_y"], marker='o', linestyle='none', color='red', ms=17,
                 label='Verworfen')
        plt.plot(df_nr["Treffer_x"], df_nr["Treffer_y"], marker='o', linestyle='none', color='green', ms=15,
                 label='Treffer')

        # set the title and labels
        plt.title(
                f"Wurfanalyse für {name_dict[nr]}: Würfe =  {würfe}, Teffer = {treffer}, Quote = {quote}%\n \
                Wenn unklar als 0/0 eingetragen. Gebt bescheid, falls was fehlt!"
                )
        plt.xlabel('Torbreite 3m = 1-15')
        plt.ylabel('Torhöhe 2m = 1-10')
        plt.legend()

        # always show entire goal (15x10)
        plt.xlim(0, 16)
        plt.ylim(0, 11)

        # save image
        img_path = os.path.join(output_dir, f"plot_nr_{nr}.png")
        plt.savefig(img_path)
        plt.close()

        # add to ppt
        slide_layout = ppt.slide_layouts[5]
        slide = ppt.slides.add_slide(slide_layout)
        left = Inches(1)
        top = Inches(1)
        slide.shapes.add_picture(img_path, left, top, width=Inches(8), height=Inches(5))

    # keeper analyse
    df_new_tw = pd.DataFrame(
        columns=['Nr.', 'Treffer_min', 'Treffer_x', 'Treffer_y', 'Verworfen_min', 'Verworfen_x', 'Verworfen_y'])
    count = 0
    for i in range(0, len(df_tw)):
        nr = df_tw.loc[i, 'Nr.']
        # split relevant columns at ;
        # FIXME: Use format_df()
        for column in df_tw.columns[3:]:
            if df_tw.loc[i, column][-1:] == ';':
                df_tw.loc[i, column] = df_tw.loc[i, column][:-1]
            entry_list = df_tw.loc[i, column].split(';')
            # df_nr = df_new
            for entry in entry_list:
                df_new_tw.loc[count, 'Nr.'] = nr
                df_new_tw.loc[count, f'{column}_min'] = entry.split(',')[0]
                if column in ['Treffer', 'Verworfen']:
                    df_new_tw.loc[count, f'{column}_y'] = entry.split(',')[2]
                    df_new_tw.loc[count, f'{column}_x'] = entry.split(',')[1]
                if column == 'Technisches Foul':
                    df_new_tw.loc[count, f'{column}_position'] = entry.split(',')[1]
                count = count + 1

        df_new_tw = df_new_tw.fillna(0)

        spieler = name_dict[nr]

        # calculate quote per position
        df_nr = df_new_tw[(df_new_tw['Nr.'] == int(nr))]
        for column in df_nr.columns[2:]:
            if 'min' not in column:
                df_nr[column] = df_nr[column].astype(int)
        treffer = len(df_nr[df_nr['Treffer_min'] != 0])
        gehalten = len(df_nr[df_nr['Verworfen_min'] != 0])
        tw_quote = round((gehalten / (treffer + gehalten) * 100), 2)

        # plot keeper analysis all positions
        plt.plot(df_nr["Verworfen_x"], df_nr["Verworfen_y"], marker='o', linestyle='none', color='green', ms=20)
        plt.plot(df_nr["Treffer_x"], df_nr["Treffer_y"], marker='o', linestyle='none', color='red', ms=15)
        plt.title(f"Torwart analyse {spieler}: rot = Treffer, grün = gehalten // Quote = {tw_quote} %")
        plt.xlabel('Torbreite 3m = 1-15')
        plt.ylabel('Torhöhe 2m = 1-10')

        # always show entire goal (15x10)
        plt.xlim(0, 16)
        plt.ylim(0, 11)

        # save image
        img_path = os.path.join(output_dir, f"plot_tw_nr_{nr}.png")
        plt.savefig(img_path)
        plt.close()

        # add to ppt
        slide_layout = ppt.slide_layouts[5]
        slide = ppt.slides.add_slide(slide_layout)
        left = Inches(1)
        top = Inches(1)
        slide.shapes.add_picture(img_path, left, top, width=Inches(8), height=Inches(5))

        # create plot per keeper for each position

        position_list = list(
            pd.DataFrame(list(df_nr['Treffer_min'].unique()) + list(df_nr['Verworfen_min'].unique()))[0].unique())

        for position in position_list:
            if position != 0:
                x = df_nr[(df_nr['Treffer_min'] == position) | (df_nr['Verworfen_min'] == position)]
                treffer = len(x[x['Treffer_min'] != 0])
                gehalten = len(x[x['Verworfen_min'] != 0])
                tw_quote = round((gehalten / (treffer + gehalten) * 100), 2)
                plt.plot(x["Verworfen_x"], x["Verworfen_y"], marker='o', linestyle='none', color='green', ms=20)
                plt.plot(x["Treffer_x"], x["Treffer_y"], marker='o', linestyle='none', color='red', ms=15)
                plt.title(f"Torwart Analyse {spieler} von {position}: // Quote = {tw_quote} %")
                plt.xlabel('Torbreite 3m = 1-15')
                plt.ylabel('Torhöhe 2m = 1-10')
                # always show entire goal (15x10)
                plt.xlim(0, 16)
                plt.ylim(0, 11)

                # save image
                img_path = os.path.join(output_dir, f"plot_tw_nr_{nr}_{position}.png")
                plt.savefig(img_path)
                plt.close()

                # add to ppt
                slide_layout = ppt.slide_layouts[5]
                slide = ppt.slides.add_slide(slide_layout)
                left = Inches(1)
                top = Inches(1)
                slide.shapes.add_picture(img_path, left, top, width=Inches(8), height=Inches(5))


    # export all to ppt and shut ppt down afterwards
    ppt_file_path = rf"{output_dir}\{filename}.ppt"
    pdf_file_path = rf"{output_dir}\{filename}.pdf"
    try:
        ppt.save(ppt_file_path)
    except PermissionError:
        print("Kollege! Mach die PowerPoint zu, es zieht! Dann nochmal probieren bitte (-_-)")
        os.system("TASKKILL /F /IM powerpnt.exe")
    try:
        PPTtoPDF(ppt_file_path, pdf_file_path)
        print('PDF wurde erstellt.')
    except:
        print('PDF muss noch erstellt werden.')


def parse_data(path: str):
    raw_data = read_data(path)

    parsed_data = pd.DataFrame(
        columns=['Nr.', 'Treffer_min', 'Treffer_x', 'Treffer_y', 'Verworfen_min', 'Verworfen_x', 'Verworfen_y',
                 'Technisches Foul_min', 'Technisches Foul_position'])

    # define keeper and non-keeper df
    nr_position_dict = {}
    for i in range(0, len(raw_data)):
        nr_position_dict[raw_data.loc[i, 'Nr.']] = raw_data.loc[i, 'Position']
    keeper_nr = []
    for nr in nr_position_dict:
        if nr_position_dict[nr] == 'TW':
            keeper_nr.append(nr)
    df_feld = raw_data[raw_data['Position'] != 'TW']
    df_tw = raw_data[raw_data['Position'] == 'TW'].reset_index(drop=True)
    parsed_data = format_df(df_feld, parsed_data)
    for column in parsed_data.columns[1:]:
        if 'position' not in column:
            parsed_data[column] = parsed_data[column].astype(int)

    return df_tw, keeper_nr, nr, parsed_data, raw_data


def read_data(path: str) -> DataFrame:
    try:
        raw_data = pd.read_excel(path)
    except PermissionError:
        print("Eh du Dulli, mach mal erst die Excel zu und probier dann nochmal.")
        sys.exit(1)
    raw_data = raw_data[['Spieler', 'Position', 'Nr.', 'Treffer', 'Verworfen', 'Technisches Foul', 'Fehlpass']].dropna(
        how='all')
    raw_data = raw_data.fillna('0,0,0')
    return raw_data


main()
