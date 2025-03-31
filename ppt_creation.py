from pptx import Presentation
from pptx.util import Inches

from Analysis import Analysis

def add_image_to_ppt(analysis: Analysis, ppt: Presentation):
    slide_layout = ppt.slide_layouts[5]
    slide = ppt.slides.add_slide(slide_layout)
    slide.shapes.add_picture(analysis.image_path, left=Inches(analysis.left), top=Inches(analysis.top),
                             width=Inches(analysis.width), height=Inches(analysis.height))


def create_ppt(analysis: Analysis):
    ppt = Presentation()

    if analysis.image_path:
        add_image_to_ppt(analysis, ppt)

    # TODO add table of contents

    add_analyse_to_ppt(analysis.sub_analyses, ppt)

    return ppt


def add_analyse_to_ppt(analyses: list[Analysis], ppt: Presentation):
    for analysis in analyses:
        if analysis.image_path:
            add_image_to_ppt(analysis, ppt)
        if len(analysis.sub_analyses) > 0:
            add_analyse_to_ppt(analysis.sub_analyses, ppt)
